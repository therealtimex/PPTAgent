import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial

from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.util import Pt

from llms import Role
from presentation import Picture, SlidePage
from utils import runs_merge


@dataclass
class HistoryMark:
    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"


class CodeExecutor:

    def __init__(self, retry_times: int):
        self.api_history = []
        self.code_history = []
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z]+\(.+\)$")

    def get_apis_docs(self, funcs: list[callable], show_example: bool = True):
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name == "slide":
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if not show_example:
                api_doc.append(signature)
                continue
            doc = inspect.getdoc(func)
            api_doc.append(f"{signature}\n\t{doc}")
        return "\n".join(api_doc)

    def execute_actions(self, actions: str, edit_slide: SlidePage):
        found_code = False
        api_calls = actions.strip().split("\n")
        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise ValueError(
                        "No code block found in the output, wrap your code with ```python```"
                    )
                if line.startswith("def"):
                    raise ValueError("The function definition should not be output.")
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise ValueError(f"The function {func} is not defined.")
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except:
                trace_msg = traceback.format_exc()
                if found_code:
                    self.code_history[-1][-1] = trace_msg
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx:])
                )
                return api_lines, trace_msg
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT


# supporting functions
def element_index(slide: SlidePage, element_id: int):
    for shape in slide:
        if shape.shape_idx == element_id:
            return shape
    raise ValueError(f"Cannot find element {element_id}, is it deleted or not exist?")


def replace_run(paragraph_id: int, run_id: int, new_text: str, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    run = runs_merge(para)[run_id]
    run.text = new_text


def clone_para(paragraph_id: int, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))


def del_run(paragraph_id: int, run_id: int, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    run = para.runs[run_id]
    run._r.getparent().remove(run._r)
    if len(para.runs) == 0:
        para._element.getparent().remove(para._element)


def set_font(
    bold: bool,
    italic: bool,
    underline: bool,
    font_size: int,
    font_color: str,
    para_id: int,
    span_id: int,
    text_shape: BaseShape,
):
    run = runs_merge(text_shape.text_frame.paragraphs[para_id])[span_id]
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if underline is not None:
        run.font.underline = underline
    if font_size is not None:
        run.font.size = Pt(font_size)
    if font_color is not None:
        run.font.color.rgb = RGBColor.from_string(font_color)


def set_size(width: int, height: int, left: int, top: int, shape: BaseShape):
    if width is not None:
        shape.width = Pt(width)
    if height is not None:
        shape.height = Pt(height)
    if left is not None:
        shape.left = Pt(left)
    if top is not None:
        shape.top = Pt(top)


# api functions
def del_span(slide: SlidePage, div_id: int, paragraph_id: int, span_id: int):
    shape = element_index(slide, div_id)
    shape._closures["delete"].append(partial(del_run, paragraph_id, span_id))


def del_image(slide: SlidePage, figure_id: int):
    shape = element_index(slide, figure_id)
    if not isinstance(shape, Picture):
        raise ValueError("The element is not a Picture.")
    slide.shapes.remove(shape)


def replace_text(
    slide: SlidePage, div_id: int, paragraph_id: int, span_id: int, text: str
):
    shape = element_index(slide, div_id)
    shape._closures["replace"].append(partial(replace_run, paragraph_id, span_id, text))


def replace_image(slide: SlidePage, figure_id: int, image_path: str):
    if not os.path.exists(image_path):
        raise ValueError(f"The image {image_path} does not exist.")
    shape = element_index(slide, figure_id)
    if not isinstance(shape, Picture):
        raise ValueError("The element is not a Picture.")
    shape.img_path = image_path


def clone_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    # The cloned paragraph will have a paragraph_id one greater than the current maximum in the parent element.
    shape = element_index(slide, div_id)
    shape._closures["clone"].append(partial(clone_para, paragraph_id))


class API_TYPES(Enum):
    Agent = [
        del_span,
        del_image,
        clone_paragraph,
        replace_text,
        replace_image,
    ]

    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs


# 加一个switch position 或者 insert? 或者split text into run, prev and next?
